# Import required modules and libraries
import re  # Regular expressions module
from langchain_community.document_loaders import PyPDFLoader  # For loading PDF files via LangChain
from docx import Document  # To work with .docx files (Word documents)
from pptx import Presentation  # To work with .pptx files (PowerPoint presentations)
from PIL import Image  # Python Imaging Library for image processing
import requests  # To make HTTP requests
from bs4 import BeautifulSoup  # For parsing HTML/XML content
from readability import Document as ReadabilityDocument  # Extracts main content from web articles
import io  # For handling byte streams
import os  # For interacting with the operating system
from dotenv import load_dotenv  # Loads environment variables from a .env file

# ================================
# --- PDF Text Extraction ---
# ================================
def extract_text_from_pdf(uploaded_file):
    # Save uploaded file locally as PDF
    with open("temp_uploaded_file.pdf", "wb") as f:
        f.write(uploaded_file.getbuffer())

    # Load the saved PDF using LangChain's PDF loader
    pdf_loader = PyPDFLoader("temp_uploaded_file.pdf")
    docs = pdf_loader.load()  # Load pages as documents

    # Patterns of text lines to exclude (e.g., file tags or empty lines)
    unwanted_patterns = [
        r'.*indd.*',  # Matches lines that contain 'indd'
        r'^\s+$'      # Matches lines that are only whitespace
    ]

    full_text = ""  # String to accumulate cleaned text

    # Iterate over each page in the document
    for page_number, doc in enumerate(docs, start=1):
        page_content = doc.page_content  # Get text content of the page
        filtered_content = []  # List to hold cleaned lines for current page

        # Filter out unwanted lines
        for line in page_content.split('\n'):
            if any(re.match(pattern, line.strip()) for pattern in unwanted_patterns):
                continue
            filtered_content.append(line.strip())

        # Add cleaned lines to full text
        for line in filtered_content:
            if line.strip():  # Skip blank lines
                full_text += line + "\n"

    return full_text  # Return cleaned and aggregated PDF text

# ================================
# --- TXT File Text Extraction ---
# ================================
def extract_text_from_txt(uploaded_file):
    # Read all lines from the text file and decode from bytes to string
    lines = uploaded_file.read().decode("utf-8").splitlines()

    # Patterns to ignore
    unwanted_patterns = [
        r'.*indd.*',
        r'^\s+$'
    ]

    full_text = ""  # Accumulate cleaned lines here
    for line in lines:
        stripped_line = line.strip()  # Remove leading/trailing whitespace
        if any(re.match(pattern, stripped_line) for pattern in unwanted_patterns):
            continue
        if stripped_line:  # Skip empty lines
            full_text += stripped_line + "\n"

    return full_text  # Return cleaned text

# ================================
# --- DOCX File Text Extraction ---
# ================================
def extract_text_from_docx(uploaded_file):
    doc = Document(uploaded_file)  # Load Word document from uploaded file

    unwanted_patterns = [
        r'.*indd.*',
        r'^\s+$'
    ]

    full_text = ""
    for para in doc.paragraphs:  # Iterate over paragraphs
        line = para.text.strip()  # Get paragraph text and trim it
        if any(re.match(pattern, line) for pattern in unwanted_patterns):
            continue
        if line:  # If not empty, add to result
            full_text += line + "\n"

    return full_text

# ================================
# --- PPTX File Text Extraction ---
# ================================
def extract_text_from_pptx(uploaded_file):
    prs = Presentation(uploaded_file)  # Load PowerPoint presentation

    unwanted_patterns = [
        r'.*indd.*',
        r'^\s+$'
    ]

    full_text = ""
    for slide_number, slide in enumerate(prs.slides, start=1):  # Iterate over slides
        for shape in slide.shapes:
            if hasattr(shape, "text"):  # Only consider shapes that have text
                for line in shape.text.split('\n'):  # Split text into lines
                    stripped_line = line.strip()
                    if any(re.match(pattern, stripped_line) for pattern in unwanted_patterns):
                        continue
                    if stripped_line:
                        full_text += stripped_line + "\n"

    return full_text

# ================================
# --- Image Text Extraction using OCR.Space API ---
# ================================
def extract_text_from_image(uploaded_file):
    api_key = os.getenv("OCR_API_KEY")  # Get OCR API key from environment

    if not api_key:
        return "Error: OCR API key is missing."  # Fail gracefully if no key

    url = 'https://api.ocr.space/parse/image'  # OCR.Space API endpoint

    # Try opening image file
    try:
        image = Image.open(uploaded_file)
    except Exception as e:
        return f"Error: Unable to open image file. {str(e)}"

    # Convert image to bytes
    img_byte_arr = io.BytesIO()
    image.save(img_byte_arr, format='PNG')  # Save image in PNG format
    img_byte_arr = img_byte_arr.getvalue()  # Get byte data

    # Send POST request to OCR.Space
    try:
        response = requests.post(
            url,
            files={'filename': ('image.png', img_byte_arr, 'image/png')},
            data={'apikey': api_key, 'language': 'eng'}
        )
    except requests.RequestException as e:
        return f"Error: Request to OCR API failed. {str(e)}"

    # Try parsing JSON from response
    try:
        result = response.json()
    except ValueError:
        return f"Error: Response from OCR.Space is not in valid JSON format. Response: {response.text}"

    # Check for OCR failure
    if result.get("IsErroredOnProcessing"):
        return f"OCR failed: {result.get('ErrorMessage', ['Unknown error'])[0]}"

    # Extract text from parsed result
    extracted_text = result["ParsedResults"][0]["ParsedText"]

    # Clean up using regex filters
    unwanted_patterns = [
        r'.*indd.*',
        r'^\s+$'
    ]

    full_text = ""
    for line in extracted_text.split('\n'):
        if any(re.match(pattern, line.strip()) for pattern in unwanted_patterns):
            continue
        if line.strip():
            full_text += line.strip() + "\n"

    return full_text

# ================================
# --- URL Article Extraction ---
# ================================
def extract_text_from_url(url):
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
        )
    }

    try:
        # Send HTTP GET request to fetch URL contents
        response = requests.get(url, headers=headers)
        response.raise_for_status()  # Raise error for bad responses

        # Use readability to extract main content
        doc = ReadabilityDocument(response.text)
        title = doc.title()  # Extract article title
        summary_html = doc.summary()  # Extract main article body as HTML

        # Parse HTML to plain text
        soup = BeautifulSoup(summary_html, 'html.parser')
        text = soup.get_text(separator='\n')  # Join blocks with newline

        # Clean with unwanted patterns
        unwanted_patterns = [
            r'.*indd.*',
            r'^\s*$'  # Empty lines
        ]

        full_text = f"Title: {title}\n\n"  # Start with the article title

        for line in text.split('\n'):
            if not any(re.match(p, line.strip()) for p in unwanted_patterns):
                full_text += line.strip() + "\n"  # Add only non-empty lines

        return full_text  # Return full extracted and cleaned article text

    except Exception as e:
        return f"Error extracting article: {e}"  # Gracefully handle all exceptions
